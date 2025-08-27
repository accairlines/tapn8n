import logging
import re
from typing import List, Dict, Any, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger(__name__)


class PPTBuilder:
    """Builds PowerPoint presentations from email content."""
    
    def __init__(self):
        self.default_slide_width = Inches(13.33)
        self.default_slide_height = Inches(7.5)
        self.title_font_size = Pt(44)
        self.subtitle_font_size = Pt(32)
        self.body_font_size = Pt(18)
        self.bullet_font_size = Pt(16)
    
    def generate_content_from_rag(self, rag_result: Dict[str, Any], slide_count: int = 5, include_charts: bool = True) -> Dict[str, Any]:
        """Generate PPT content structure from RAG results."""
        try:
            # Extract key information
            answer = rag_result.get('answer', '')
            sources = rag_result.get('sources', [])
            
            # Parse the answer to extract structured content
            content_structure = self._parse_llama_response(answer)
            
            # If parsing failed, create a simple structure
            if not content_structure:
                content_structure = self._create_fallback_structure(answer, sources, slide_count)
            
            # Add metadata
            content_structure['metadata'] = {
                'total_slides': len(content_structure.get('slides', [])),
                'include_charts': include_charts,
                'sources_count': len(sources),
                'generation_timestamp': rag_result.get('timestamp', '')
            }
            
            return content_structure
            
        except Exception as e:
            logger.error(f"Failed to generate PPT content: {e}")
            return self._create_error_structure(str(e))
    
    def build_pptx(self, content_structure: Dict[str, Any]) -> bytes:
        """Build PPTX file from content structure."""
        try:
            # Create presentation
            prs = Presentation()
            
            # Set slide dimensions
            prs.slide_width = self.default_slide_width
            prs.slide_height = self.default_slide_height
            
            # Add title slide
            self._add_title_slide(prs, content_structure)
            
            # Add content slides
            slides = content_structure.get('slides', [])
            for slide_data in slides:
                self._add_content_slide(prs, slide_data)
            
            # Add sources slide if available
            sources = content_structure.get('sources', [])
            if sources:
                self._add_sources_slide(prs, sources)
            
            # Save to bytes
            import io
            pptx_buffer = io.BytesIO()
            prs.save(pptx_buffer)
            pptx_buffer.seek(0)
            
            return pptx_buffer.getvalue()
            
        except Exception as e:
            logger.error(f"Failed to build PPTX: {e}")
            raise
    
    def _parse_llama_response(self, response: str) -> Dict[str, Any]:
        """Parse LLaMA response to extract structured content."""
        try:
            lines = response.split('\n')
            content_structure = {
                'title': '',
                'slides': [],
                'sources': []
            }
            
            current_slide = None
            in_slide_content = False
            
            for line in lines:
                line = line.strip()
                if not line:
                    continue
                
                # Extract title
                if 'título' in line.lower() and ':' in line:
                    content_structure['title'] = line.split(':', 1)[1].strip()
                
                # Extract slide information
                elif 'slide' in line.lower() or 'slide' in line.lower():
                    if current_slide:
                        content_structure['slides'].append(current_slide)
                    
                    current_slide = {
                        'title': line,
                        'content': [],
                        'bullet_points': []
                    }
                    in_slide_content = True
                
                # Extract slide content
                elif current_slide and in_slide_content:
                    if line.startswith('-') or line.startswith('•'):
                        current_slide['bullet_points'].append(line[1:].strip())
                    else:
                        current_slide['content'].append(line)
            
            # Add the last slide
            if current_slide:
                content_structure['slides'].append(current_slide)
            
            return content_structure
            
        except Exception as e:
            logger.error(f"Failed to parse LLaMA response: {e}")
            return {}
    
    def _create_fallback_structure(self, answer: str, sources: List[Dict[str, Any]], slide_count: int) -> Dict[str, Any]:
        """Create fallback content structure when parsing fails."""
        content_structure = {
            'title': 'Análise de Emails',
            'slides': [],
            'sources': sources
        }
        
        # Split answer into sentences for slides
        sentences = re.split(r'[.!?]+', answer)
        sentences = [s.strip() for s in sentences if s.strip()]
        
        # Create slides from sentences
        slides_per_slide = max(1, len(sentences) // slide_count)
        
        for i in range(0, len(sentences), slides_per_slide):
            slide_sentences = sentences[i:i + slides_per_slide]
            slide_data = {
                'title': f'Slide {len(content_structure["slides"]) + 1}',
                'content': slide_sentences,
                'bullet_points': slide_sentences
            }
            content_structure['slides'].append(slide_data)
            
            if len(content_structure['slides']) >= slide_count:
                break
        
        return content_structure
    
    def _create_error_structure(self, error_message: str) -> Dict[str, Any]:
        """Create error content structure."""
        return {
            'title': 'Erro na Geração',
            'slides': [{
                'title': 'Erro',
                'content': [error_message],
                'bullet_points': [error_message]
            }],
            'sources': []
        }
    
    def _add_title_slide(self, prs: Presentation, content_structure: Dict[str, Any]):
        """Add title slide to presentation."""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = content_structure.get('title', 'Análise de Emails')
        
        # Set subtitle
        subtitle = slide.placeholders[1]
        subtitle.text = f"Gerado automaticamente com base nos emails indexados\n{content_structure.get('metadata', {}).get('generation_timestamp', '')}"
        
        # Format title
        title.text_frame.paragraphs[0].font.size = self.title_font_size
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)
    
    def _add_content_slide(self, prs: Presentation, slide_data: Dict[str, Any]):
        """Add content slide to presentation."""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = slide_data.get('title', 'Slide sem título')
        title.text_frame.paragraphs[0].font.size = self.subtitle_font_size
        title.text_frame.paragraphs[0].font.bold = True
        
        # Set content
        content = slide.placeholders[1]
        text_frame = content.text_frame
        
        # Clear default text
        text_frame.clear()
        
        # Add bullet points
        bullet_points = slide_data.get('bullet_points', [])
        if bullet_points:
            for i, point in enumerate(bullet_points):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = point
                p.font.size = self.bullet_font_size
                p.level = 0
        
        # Add regular content if no bullet points
        else:
            content_list = slide_data.get('content', [])
            for i, line in enumerate(content_list):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = line
                p.font.size = self.body_font_size
    
    def _add_sources_slide(self, prs: Presentation, sources: List[Dict[str, Any]]):
        """Add sources slide to presentation."""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = "Fontes dos Emails"
        title.text_frame.paragraphs[0].font.size = self.subtitle_font_size
        title.text_frame.paragraphs[0].font.bold = True
        
        # Set content
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        # Add sources
        for i, source in enumerate(sources):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            subject = source.get('subject', 'Assunto desconhecido')
            date = source.get('date', 'Data desconhecida')
            file_path = source.get('file', 'Ficheiro desconhecido')
            
            p.text = f"• {subject} | {date} | {file_path}"
            p.font.size = self.bullet_font_size
            p.level = 0
    
    def _add_chart_slide(self, prs: Presentation, chart_data: Dict[str, Any]):
        """Add chart slide to presentation (placeholder for future enhancement)."""
        # This is a placeholder for future chart generation capabilities
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Gráfico de Dados"
        
        content = slide.placeholders[1]
        content.text = "Funcionalidade de gráficos será implementada em versões futuras."
